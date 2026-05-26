from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "0525_pre_methods_overview.pptx"

FONT = "Noto Sans CJK SC"


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_run_font(run, size: float, color: str, *, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 14,
    color: str = "#17202A",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size, color, bold=bold)


def add_round_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    *,
    fill: str,
    line: str,
    title_color: str = "#17202A",
    body_color: str = "#243342",
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_font(run, 11.5, title_color, bold=True)

    if body:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = body
        set_run_font(run, 8.4, body_color)


def add_header_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    subtitle: str,
    *,
    fill: str,
    accent: str,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(accent)
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_font(run, 15.2, "#102027", bold=True)

    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = subtitle
    set_run_font(run, 8.7, "#41505D")


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    bullets: list[str],
    *,
    color: str = "#263238",
    size: float = 8.8,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(1.0)


def add_lane(slide, x: float, y: float, w: float, h: float, fill: str, line: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)


def add_arrow(slide, x: float, y: float, *, color: str, size: float = 14) -> None:
    add_text(slide, x, y, 0.16, 0.18, "→", size=size, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_ram_to_residual_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#F7F8F4")

    add_text(
        slide,
        0.46,
        0.22,
        12.45,
        0.34,
        "为什么从 RAM→3ch→DAV2 转向 Residual Correction",
        size=19.4,
        color="#102027",
        bold=True,
    )
    add_text(
        slide,
        0.47,
        0.58,
        12.2,
        0.28,
        "旧路线是 input replacement：让 RAW-like/RAM 输出伪装成 RGB 进入 DAV2；新路线是 output correction：保留 RGB-DAV2 主先验，只在失败区域做受控修正。",
        size=9.4,
        color="#53646F",
    )

    add_lane(slide, 0.42, 1.0, 5.95, 5.05, "#F1E8E5", "#C88372")
    add_lane(slide, 6.96, 1.0, 5.95, 5.05, "#E7F0EA", "#79A887")

    add_header_box(
        slide,
        0.66,
        1.18,
        5.48,
        0.62,
        "旧路线：把 RAM 输出当成 DAV2 输入",
        "raw4 -> base_rgb -> RamCore3 -> x3 -> frozen DAV2 -> depth",
        fill="#F8D8CF",
        accent="#B56654",
    )
    add_header_box(
        slide,
        7.2,
        1.18,
        5.48,
        0.62,
        "当前路线：保留 DAV2，只预测修正量",
        "RGB -> frozen DAV2 -> D0；RAW/RAM 特征只进入 gate + delta head",
        fill="#D9EBDD",
        accent="#5A9A68",
    )

    old_xs = [0.75, 1.92, 3.09, 4.26]
    old_steps = [
        ("raw4", "packed\nBayer"),
        ("RAM", "learned\n3ch x3"),
        ("DAV2", "expects\nRGB stats"),
        ("Depth", "full map\nrewritten"),
    ]
    for x, (title, body) in zip(old_xs, old_steps):
        add_round_box(slide, x, 2.06, 1.02, 0.72, title, body, fill="#FFFFFF", line="#CF9B91")
    for x in [1.78, 2.95, 4.12]:
        add_arrow(slide, x, 2.31, color="#A64D3D")

    new_top_xs = [7.28, 8.56, 9.84]
    new_top_steps = [
        ("RGB", "original\nimage"),
        ("DAV2", "frozen\nmain path"),
        ("D0", "base depth\nprior"),
    ]
    for x, (title, body) in zip(new_top_xs, new_top_steps):
        add_round_box(slide, x, 1.98, 1.08, 0.68, title, body, fill="#FFFFFF", line="#94BE9B")
    for x in [8.38, 9.66]:
        add_arrow(slide, x, 2.21, color="#4C8B58")

    new_bottom_xs = [7.28, 8.56, 9.84, 11.12]
    new_bottom_steps = [
        ("raw4", "packed\nBayer"),
        ("RAM", "x3 /\nffm_mid"),
        ("Head", "gate +\ndelta"),
        ("Final", "D0 +\ng*delta"),
    ]
    for x, (title, body) in zip(new_bottom_xs, new_bottom_steps):
        add_round_box(slide, x, 3.02, 1.08, 0.68, title, body, fill="#FFFFFF", line="#94BE9B")
    for x in [8.38, 9.66, 10.94]:
        add_arrow(slide, x, 3.25, color="#4C8B58")

    add_bullets(
        slide,
        0.72,
        3.18,
        5.28,
        1.5,
        [
            "输入分布失配：DAV2 是 RGB depth foundation model，x3/RAM 输出不是自然 RGB",
            "会重写 DAV2 依赖的语义-几何映射：shading、texture、sky/far priors、layout",
            "RAM 的成功证据主要来自 detection；dense depth 对整图一致性和边界更敏感",
            "synthetic RAW-like 由 sRGB 反推，不能当成真实新增 sensor information",
        ],
        color="#57332D",
        size=7.9,
    )

    add_bullets(
        slide,
        7.25,
        4.0,
        5.28,
        1.24,
        [
            "RGB 仍走 frozen DAV2：保住强 prior，避免破坏输入统计",
            "RAW-like/RAM branch 只回答 where to fix 与 how much to fix",
            "zero-init delta、gate bias<0、gate/residual regularization 让初始行为接近 D0",
            "更容易证明：RAW-like 是 complementary residual cue，而不是主输入替代物",
        ],
        color="#2B4932",
        size=7.9,
    )

    add_round_box(
        slide,
        0.78,
        4.92,
        5.12,
        0.72,
        "旧路线的问题定位",
        "不是 RAM 一定无效，而是它把辅助表征放在 DAV2 前面，导致基础模型必须在非 RGB 统计上重新解释整张图。",
        fill="#F8D8CF",
        line="#B56654",
        title_color="#4B2017",
        body_color="#51312B",
    )
    add_round_box(
        slide,
        7.3,
        5.38,
        5.1,
        0.72,
        "新路线的关键分工",
        "DAV2 负责整图形状和语义先验；RAW-like/RAM 特征只作为局部误差校正信号，最终 pred = D0_norm + gate * delta。",
        fill="#D9EBDD",
        line="#5A9A68",
        title_color="#173B21",
        body_color="#274A30",
    )

    add_round_box(
        slide,
        0.55,
        6.2,
        12.25,
        0.62,
        "一句话叙事",
        "从“让 RAW-like 替代 RGB 进入 DAV2”，改为“让 RAW-like 在不破坏 DAV2 RGB prior 的前提下修正 DAV2 的 failure regions”。",
        fill="#FFFFFF",
        line="#B8C2C8",
        title_color="#102027",
        body_color="#2D3A42",
    )

    add_text(
        slide,
        0.47,
        7.05,
        12.42,
        0.2,
        "Sources: plans/deep-research-report.md; plans/0524_new/RAMCore3_DAV2_RAW_residual_vkitti_mseries_execution_plan.md; foundation/engine/models/raw_residual_dav2.py",
        size=6.5,
        color="#75838C",
    )


def add_network_training_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#F7F8F4")

    add_text(
        slide,
        0.42,
        0.2,
        12.55,
        0.34,
        "当前网络训练流程：RGB D0 prior + RAW/RAM residual correction",
        size=18.3,
        color="#102027",
        bold=True,
    )
    add_text(
        slide,
        0.43,
        0.56,
        12.15,
        0.28,
        "同一个样本同时走 RGB 主路、RAW 分支和 GT target 分支；DAV2 frozen/no_grad，训练 RAM + ResidualGateHead，最终只学习对 D0_norm 的 gate * delta 修正。",
        size=8.9,
        color="#53646F",
    )

    add_lane(slide, 0.42, 1.0, 3.32, 4.72, "#E6F0EC", "#78A58D")
    add_lane(slide, 4.02, 1.0, 5.18, 4.72, "#F3EAE1", "#C49A6C")
    add_lane(slide, 9.55, 1.0, 3.36, 4.72, "#ECE9F3", "#9D91B6")

    add_header_box(
        slide,
        0.66,
        1.2,
        2.84,
        0.58,
        "Batch sample",
        "dataset 一次 crop/flip 后产生四类张量",
        fill="#D9EBE1",
        accent="#5F9A79",
    )
    add_round_box(
        slide,
        0.82,
        2.04,
        2.52,
        0.58,
        "image",
        "[3,H,W] ImageNet norm；给 frozen DAV2",
        fill="#FFFFFF",
        line="#8DB79D",
    )
    add_round_box(
        slide,
        0.82,
        2.88,
        2.52,
        0.58,
        "raw",
        "[4,H,W] packed Bayer；sensor_linear_dual",
        fill="#FFFFFF",
        line="#8DB79D",
    )
    add_round_box(
        slide,
        0.82,
        3.72,
        2.52,
        0.58,
        "depth",
        "metric GT depth；用于 inverse-depth target",
        fill="#FFFFFF",
        line="#8DB79D",
    )
    add_round_box(
        slide,
        0.82,
        4.56,
        2.52,
        0.58,
        "valid_mask",
        "过滤无效/超 min-max depth 像素",
        fill="#FFFFFF",
        line="#8DB79D",
    )

    add_header_box(
        slide,
        4.26,
        1.2,
        4.7,
        0.58,
        "Forward graph",
        "RGB 路给 D0；RAW/RAM 路给 residual feature；GT 路给 y_norm",
        fill="#F2DDC9",
        accent="#B9844F",
    )
    add_round_box(slide, 4.35, 2.0, 1.12, 0.5, "image", "RGB norm", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 5.55, 2.15, color="#A36F37", size=12)
    add_round_box(slide, 5.72, 2.0, 1.16, 0.5, "DAV2", "frozen\nno_grad", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 6.96, 2.15, color="#A36F37", size=12)
    add_round_box(slide, 7.13, 2.0, 0.9, 0.5, "D0", "detach", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 8.12, 2.15, color="#A36F37", size=12)
    add_round_box(slide, 8.28, 2.0, 0.66, 0.5, "D0n", "norm", fill="#FFFFFF", line="#D0A271")

    add_round_box(slide, 4.35, 3.0, 0.92, 0.5, "raw4", "packed", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 5.36, 3.15, color="#A36F37", size=12)
    add_round_box(slide, 5.52, 3.0, 1.24, 0.5, "base_rgb", "R/G/B", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 6.84, 3.15, color="#A36F37", size=12)
    add_round_box(slide, 7.0, 3.0, 0.96, 0.5, "RAM", "trainable", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 8.04, 3.15, color="#A36F37", size=12)
    add_round_box(slide, 8.2, 3.0, 0.74, 0.5, "feat", "ffm64", fill="#FFFFFF", line="#D0A271")

    add_round_box(slide, 4.35, 4.02, 1.16, 0.5, "depth", "metric", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 5.6, 4.17, color="#A36F37", size=12)
    add_round_box(slide, 5.76, 4.02, 1.1, 0.5, "inv_gt", "1/depth", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 6.94, 4.17, color="#A36F37", size=12)
    add_round_box(slide, 7.1, 4.02, 1.16, 0.5, "y_norm", "target", fill="#FFFFFF", line="#D0A271")

    add_round_box(
        slide,
        4.78,
        5.02,
        1.72,
        0.48,
        "concat",
        "D0_norm + residual feature",
        fill="#FFF7EE",
        line="#C9925F",
        title_color="#5A3110",
        body_color="#5A3110",
    )
    add_arrow(slide, 6.58, 5.18, color="#A36F37", size=12)
    add_round_box(
        slide,
        6.78,
        5.02,
        1.22,
        0.48,
        "Head",
        "gate, delta",
        fill="#FFF7EE",
        line="#C9925F",
        title_color="#5A3110",
        body_color="#5A3110",
    )
    add_arrow(slide, 8.06, 5.18, color="#A36F37", size=12)
    add_round_box(slide, 8.18, 5.02, 0.78, 0.48, "pred", "D0+g*d", fill="#FFF7EE", line="#C9925F")

    add_header_box(
        slide,
        9.78,
        1.2,
        2.88,
        0.58,
        "Loss set",
        "pred 与 y_norm 对齐，同时限制修正范围",
        fill="#E4DEF0",
        accent="#8F7AB2",
    )
    add_bullets(
        slide,
        9.9,
        2.08,
        2.58,
        2.0,
        [
            "L_depth: pred vs y_norm",
            "L_grad: depth gradient shape",
            "L_keep: D0-good regions 少改",
            "L_res: 修正幅度小",
            "L_gate: gate 稀疏",
            "L_gate_sup: gate 对齐 D0 error mask",
        ],
        color="#45365B",
        size=7.6,
    )
    add_round_box(
        slide,
        9.9,
        4.34,
        2.58,
        0.74,
        "总 loss",
        "L = L_depth + 0.5L_grad + 0.1L_keep + 0.01L_res + 0.005L_gate + 0.05L_gate_sup",
        fill="#F6F2FB",
        line="#9C8AB9",
        title_color="#3E2C55",
        body_color="#493A5E",
    )

    add_round_box(
        slide,
        0.55,
        6.02,
        12.25,
        0.62,
        "三组对照的 head_input 差异",
        "M2 = concat(D0_norm, ffm_mid)；M1 = concat(D0_norm, x3)；M3 = concat(D0_norm, x3, ffm_mid)。C1 = concat(D0_norm, RGB)，C2 = D0_norm only。",
        fill="#FFFFFF",
        line="#B8C2C8",
        title_color="#102027",
        body_color="#2D3A42",
    )
    add_text(
        slide,
        0.46,
        7.05,
        12.45,
        0.2,
        "Sources: foundation/tools/train_vkitti2_raw_residual.py; foundation/engine/models/raw_residual_dav2.py; foundation/engine/datasets/vkitti2_raw.py; plans/0525_m2c1c2/README.md",
        size=6.1,
        color="#75838C",
    )


def add_loss_explanation_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#F7F8F4")

    add_text(
        slide,
        0.42,
        0.2,
        12.55,
        0.34,
        "Loss 设计：让 residual 改该改的地方，不破坏 D0 prior",
        size=18.8,
        color="#102027",
        bold=True,
    )
    add_text(
        slide,
        0.43,
        0.58,
        12.0,
        0.28,
        "target 在 normalized inverse-depth 空间；gate 的监督来自 D0 与 GT 的误差分布，所以训练目标天然是“在 D0 failure regions 做受控 correction”。",
        size=9.0,
        color="#53646F",
    )

    add_lane(slide, 0.42, 0.98, 12.5, 1.18, "#E8EEF4", "#8AA5BB")
    add_round_box(
        slide,
        0.76,
        1.22,
        3.55,
        0.58,
        "Target",
        "y_norm = robust_normalize(1 / depth, valid_mask)",
        fill="#FFFFFF",
        line="#91A9BC",
    )
    add_round_box(
        slide,
        4.86,
        1.22,
        3.38,
        0.58,
        "D0 error mask",
        "e0 = |D0_norm - y_norm|；m_error = q80->q95 soft mask",
        fill="#FFFFFF",
        line="#91A9BC",
    )
    add_round_box(
        slide,
        8.76,
        1.22,
        3.36,
        0.58,
        "Correction",
        "gate_delta = gate * delta；pred = D0_norm + gate_delta",
        fill="#FFFFFF",
        line="#91A9BC",
    )

    boxes = [
        (
            0.62,
            2.38,
            "L_depth  weight=1.0",
            "mean |pred - y_norm| over valid pixels\n主监督：让最终预测接近 GT target。",
            "#E6F0EC",
            "#78A58D",
        ),
        (
            6.82,
            2.38,
            "L_grad  weight=0.5",
            "gradient_l1(pred, y_norm)\n约束局部斜率和边界，减少只对齐均值但形状错误。",
            "#E6F0EC",
            "#78A58D",
        ),
        (
            0.62,
            3.34,
            "L_keep  weight=0.1",
            "mean((1 - m_error) * |gate_delta|)\nD0 已经好的区域，惩罚 residual 乱改。",
            "#F3EAE1",
            "#C49A6C",
        ),
        (
            6.82,
            3.34,
            "L_res  weight=0.01",
            "mean |gate_delta|\n限制最终 correction 幅度，让模型保守地修正。",
            "#F3EAE1",
            "#C49A6C",
        ),
        (
            0.62,
            4.3,
            "L_gate  weight=0.005",
            "mean gate\n让 gate 不要整图打开，保持 correction 稀疏。",
            "#ECE9F3",
            "#9D91B6",
        ),
        (
            6.82,
            4.3,
            "L_gate_sup  weight=0.05",
            "BCE(gate, m_error)\n用 D0 high-error mask 指导 gate：错得多的地方更该开。",
            "#ECE9F3",
            "#9D91B6",
        ),
    ]
    for x, y, title, body, fill, line in boxes:
        add_round_box(
            slide,
            x,
            y,
            5.72,
            0.72,
            title,
            body,
            fill=fill,
            line=line,
            title_color="#102027",
            body_color="#2D3A42",
        )

    add_round_box(
        slide,
        0.76,
        5.56,
        11.58,
        0.58,
        "完整目标",
        "loss = L_depth + 0.5 * L_grad + 0.1 * L_keep + 0.01 * L_res + 0.005 * L_gate + 0.05 * L_gate_sup",
        fill="#FFFFFF",
        line="#B8C2C8",
        title_color="#102027",
        body_color="#2D3A42",
    )
    add_round_box(
        slide,
        0.76,
        6.28,
        11.58,
        0.44,
        "直观理解",
        "L_depth/L_grad 负责把 pred 拉向 GT；L_keep/L_res/L_gate 负责不乱改；L_gate_sup 负责告诉 gate 应该在哪里打开。",
        fill="#FFFFFF",
        line="#B8C2C8",
        title_color="#102027",
        body_color="#2D3A42",
    )
    add_text(
        slide,
        0.46,
        7.05,
        12.45,
        0.2,
        "Sources: foundation/tools/train_vkitti2_raw_residual.py::compute_residual_loss; foundation/engine/models/raw_residual_dav2.py::RawResidualDAV2.forward",
        size=6.3,
        color="#75838C",
    )


def add_raw_d0_network_detail_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#F7F8F4")

    add_text(
        slide,
        0.42,
        0.2,
        12.55,
        0.34,
        "M2 raw + D0_norm 具体网络结构",
        size=19.0,
        color="#102027",
        bold=True,
    )
    add_text(
        slide,
        0.43,
        0.58,
        12.0,
        0.28,
        "这一页展开 residual_feature_source=ffm_mid 的路径：RAW/RAM 负责产生 64ch feature，D0_norm 提供 1ch prior，二者拼成 65ch 后进入 ResidualGateHead。",
        size=9.0,
        color="#53646F",
    )

    add_lane(slide, 0.42, 1.0, 5.96, 4.92, "#E6F0EC", "#78A58D")
    add_lane(slide, 6.78, 1.0, 6.14, 4.92, "#F3EAE1", "#C49A6C")

    add_header_box(
        slide,
        0.66,
        1.2,
        5.46,
        0.54,
        "RAW/RAM feature path",
        "trainable RamCore3；输出 x3 作为诊断，M2 head 使用 ffm_mid",
        fill="#D9EBE1",
        accent="#5F9A79",
    )
    add_round_box(slide, 0.72, 2.05, 1.08, 0.54, "raw4", "[B,4,H,W]\nR Gr Gb B", fill="#FFFFFF", line="#8DB79D")
    add_arrow(slide, 1.9, 2.22, color="#5A8D6F", size=12)
    add_round_box(
        slide,
        2.1,
        2.05,
        1.44,
        0.54,
        "base_rgb",
        "[R,(Gr+Gb)/2,B]\n[B,3,H,W]",
        fill="#FFFFFF",
        line="#8DB79D",
    )
    add_arrow(slide, 3.65, 2.22, color="#5A8D6F", size=12)
    add_round_box(slide, 3.86, 2.05, 1.16, 0.54, "RPEnc", "resize 256\nz [B,128]", fill="#FFFFFF", line="#8DB79D")

    add_round_box(
        slide,
        0.72,
        3.0,
        5.18,
        0.7,
        "4 parallel ISP branches",
        "WB gain -> 3ch；CCM 3x3 -> 3ch；Gamma -> 3ch；Brightness -> 3ch。每个 branch 都由 z 解码参数，并作用在 base_rgb 上。",
        fill="#FFFFFF",
        line="#8DB79D",
        title_color="#173B21",
        body_color="#2F5140",
    )
    add_round_box(slide, 0.72, 4.05, 1.28, 0.52, "concat", "x_cat\n[B,12,H,W]", fill="#FFFFFF", line="#8DB79D")
    add_arrow(slide, 2.12, 4.2, color="#5A8D6F", size=12)
    add_round_box(
        slide,
        2.32,
        4.05,
        1.72,
        0.52,
        "FFM3",
        "12 -> 16 -> 64 -> 16 -> 3",
        fill="#FFFFFF",
        line="#8DB79D",
    )
    add_arrow(slide, 4.16, 4.2, color="#5A8D6F", size=12)
    add_round_box(slide, 4.36, 3.86, 1.26, 0.44, "ffm_mid", "[B,64,H,W]", fill="#DFF1E7", line="#5F9A79")
    add_round_box(slide, 4.36, 4.48, 1.26, 0.44, "x3", "[B,3,H,W]", fill="#FFFFFF", line="#8DB79D")

    add_header_box(
        slide,
        7.02,
        1.2,
        5.66,
        0.54,
        "D0_norm + ResidualGateHead",
        "U-Net-like head；输出 delta 与 gate，最终 pred = D0_norm + gate * delta",
        fill="#F2DDC9",
        accent="#B9844F",
    )
    add_round_box(slide, 7.0, 2.05, 1.26, 0.52, "D0_norm", "[B,1,H,W]\ndetach", fill="#FFFFFF", line="#D0A271")
    add_text(slide, 8.36, 2.18, 0.18, 0.2, "+", size=15, color="#A36F37", bold=True, align=PP_ALIGN.CENTER)
    add_round_box(slide, 8.62, 2.05, 1.28, 0.52, "ffm_mid", "[B,64,H,W]", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 10.02, 2.22, color="#A36F37", size=12)
    add_round_box(slide, 10.24, 2.05, 1.4, 0.52, "head_input", "[B,65,H,W]", fill="#FFF7EE", line="#C9925F")

    head_y = 3.02
    head_steps = [
        ("stem", "Conv3x3\n65->64"),
        ("enc0", "ResBlock\n64"),
        ("down1", "stride2\n64->128"),
        ("down2", "stride2\n128->256"),
    ]
    head_xs = [7.0, 8.25, 9.5, 10.75]
    for x, (title, body) in zip(head_xs, head_steps):
        add_round_box(slide, x, head_y, 1.02, 0.56, title, body, fill="#FFFFFF", line="#D0A271")
    for x in [8.05, 9.3, 10.55]:
        add_arrow(slide, x, head_y + 0.18, color="#A36F37", size=12)

    add_round_box(slide, 8.0, 4.06, 1.08, 0.56, "up1", "interp + skip\n256+128->128", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 9.18, 4.24, color="#A36F37", size=12)
    add_round_box(slide, 9.4, 4.06, 1.08, 0.56, "up2", "interp + skip\n128+64->64", fill="#FFFFFF", line="#D0A271")
    add_arrow(slide, 10.58, 4.24, color="#A36F37", size=12)
    add_round_box(slide, 10.82, 3.84, 1.18, 0.48, "delta", "0.5*tanh(.)", fill="#FFFFFF", line="#D0A271")
    add_round_box(slide, 10.82, 4.52, 1.18, 0.48, "gate", "sigmoid(.)", fill="#FFFFFF", line="#D0A271")
    add_round_box(
        slide,
        7.08,
        5.08,
        5.02,
        0.56,
        "initial behavior",
        "delta last conv = 0；gate bias = -4；初始 pred 接近 D0_norm。",
        fill="#FFF7EE",
        line="#C9925F",
        title_color="#5A3110",
        body_color="#5A3110",
    )

    add_round_box(
        slide,
        0.55,
        6.15,
        12.25,
        0.6,
        "训练边界",
        "Frozen/no_grad: DAV2 与 D0_norm；Trainable: RamCore3 + ResidualGateHead。当前 M2 使用 [D0_norm, ffm_mid]，不把 x3 再送进 DAV2。",
        fill="#FFFFFF",
        line="#B8C2C8",
        title_color="#102027",
        body_color="#2D3A42",
    )
    add_text(
        slide,
        0.46,
        7.05,
        12.45,
        0.2,
        "Sources: finetune_stf/models/raw_ram.py::RamCore3; foundation/engine/models/raw_residual_dav2.py::RawResidualDAV2 and ResidualGateHead",
        size=6.3,
        color="#75838C",
    )


def add_c1_c2_control_comparison_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#F7F8F4")

    add_text(
        slide,
        0.42,
        0.2,
        12.55,
        0.34,
        "正式实验 control：C1 RGB residual vs C2 D0-only residual",
        size=18.4,
        color="#102027",
        bold=True,
    )
    add_text(
        slide,
        0.43,
        0.58,
        12.1,
        0.28,
        "C1/C2 不使用 RAW/RAM，专门回答 residual gain 来自哪里：RGB appearance 是否有额外帮助，还是 D0 几何本身已经足够做 post-processing。",
        size=9.0,
        color="#53646F",
    )

    add_lane(slide, 0.42, 1.0, 12.5, 1.1, "#E8EEF4", "#8AA5BB")
    common_steps = [
        ("image", "RGB norm\n[B,3,H,W]"),
        ("DAV2", "frozen\nno_grad"),
        ("D0", "detach"),
        ("D0_norm", "robust norm\n[B,1,H,W]"),
        ("Head", "same U-Net-like\nResidualGateHead"),
        ("pred", "D0_norm +\ngate*delta"),
    ]
    xs = [0.78, 2.5, 4.22, 5.94, 8.02, 10.24]
    widths = [1.1, 1.1, 0.92, 1.34, 1.58, 1.36]
    for x, w, (title, body) in zip(xs, widths, common_steps):
        add_round_box(slide, x, 1.32, w, 0.54, title, body, fill="#FFFFFF", line="#91A9BC")
    for x in [1.96, 3.68, 5.22, 7.42, 9.72]:
        add_arrow(slide, x, 1.49, color="#5E7D95", size=12)

    add_lane(slide, 0.42, 2.38, 5.95, 3.45, "#E6F0EC", "#78A58D")
    add_lane(slide, 6.96, 2.38, 5.95, 3.45, "#ECE9F3", "#9D91B6")

    add_header_box(
        slide,
        0.66,
        2.58,
        5.48,
        0.62,
        "C1：RGB residual control",
        "head_input = concat(D0_norm, image_rgb_norm)",
        fill="#D9EBE1",
        accent="#5F9A79",
    )
    add_header_box(
        slide,
        7.2,
        2.58,
        5.48,
        0.62,
        "C2：D0-only residual control",
        "head_input = D0_norm.unsqueeze(1)",
        fill="#E4DEF0",
        accent="#8F7AB2",
    )

    add_round_box(slide, 0.78, 3.46, 1.28, 0.5, "D0_norm", "1ch", fill="#FFFFFF", line="#8DB79D")
    add_text(slide, 2.16, 3.58, 0.18, 0.2, "+", size=15, color="#5A8D6F", bold=True, align=PP_ALIGN.CENTER)
    add_round_box(slide, 2.42, 3.46, 1.28, 0.5, "RGB", "3ch", fill="#FFFFFF", line="#8DB79D")
    add_arrow(slide, 3.84, 3.6, color="#5A8D6F", size=12)
    add_round_box(slide, 4.04, 3.46, 1.34, 0.5, "head_input", "4ch", fill="#FFFFFF", line="#8DB79D")

    add_round_box(slide, 7.42, 3.46, 1.42, 0.5, "D0_norm", "1ch", fill="#FFFFFF", line="#A394C0")
    add_arrow(slide, 8.98, 3.6, color="#7B689F", size=12)
    add_round_box(slide, 9.18, 3.46, 1.42, 0.5, "head_input", "1ch", fill="#FFFFFF", line="#A394C0")

    add_bullets(
        slide,
        0.82,
        4.24,
        5.05,
        0.96,
        [
            "测试普通 RGB appearance 是否能在 D0 之外继续提供 residual cue",
            "trainable params: 2,883,586；只比 C2 多第一层 3 个输入通道",
            "风险：appearance shortcut，Scene20 holdout 上 early best 后回退",
        ],
        color="#2F5140",
        size=7.6,
    )
    add_bullets(
        slide,
        7.36,
        4.24,
        5.05,
        0.96,
        [
            "测试 D0 自身是否已足够做局部 calibration / post-processing",
            "trainable params: 2,881,858；只训练 residual head",
            "当前最强 control：M 系列必须超过 C2 才能证明 RAW/RAM cue 额外有效",
        ],
        color="#45365B",
        size=7.6,
    )

    add_round_box(
        slide,
        0.62,
        5.98,
        12.08,
        0.62,
        "当前 formal 结果速记",
        "VKITTI best abs_rel: C2 0.12100 @11 < C1 0.12574 @03；KITTI best abs_rel: C2 0.09503 @06 < C1 0.09843 @06。C2 更稳定，也定义了后续 M-series 的必须超过基线。",
        fill="#FFFFFF",
        line="#B8C2C8",
        title_color="#102027",
        body_color="#2D3A42",
    )
    add_text(
        slide,
        0.46,
        7.05,
        12.45,
        0.2,
        "Sources: foundation/engine/models/dav2_residual_control.py; foundation/tools/train_vkitti2_residual_control.py; plans/0525_m2c1c2/README.md",
        size=6.3,
        color="#75838C",
    )


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#F6F8F6")

    add_text(
        slide,
        0.42,
        0.22,
        12.55,
        0.32,
        "两种 Unprocessing 路线对比",
        size=20,
        color="#102027",
        bold=True,
    )
    add_text(
        slide,
        0.43,
        0.58,
        12.0,
        0.28,
        "共同点：都把 RGB/sRGB 转成 4-channel packed Bayer RAW-like；差异在反推 ISP 假设、随机化语义与实验可比性。",
        size=9.8,
        color="#53646F",
    )

    add_lane(slide, 0.42, 0.98, 5.95, 4.95, "#E9F2F3", "#79AAB2")
    add_lane(slide, 6.96, 0.98, 5.95, 4.95, "#FFF1E5", "#D99B5C")

    add_header_box(
        slide,
        0.66,
        1.18,
        5.48,
        0.62,
        "A. 当前训练：sensor_linear_dual",
        "Brooks-style 在线合成；VKITTI2Raw 在 dataloader 中生成 raw4",
        fill="#D8EAED",
        accent="#4D9AA6",
    )
    add_header_box(
        slide,
        7.2,
        1.18,
        5.48,
        0.62,
        "B. 0524 计划：InvISP / RAW-Adapter-style",
        "参考 RAW-Adapter：RGB -> InvISP raw-RGB -> inverse WB -> mosaic",
        fill="#FFE4CA",
        accent="#CF7A30",
    )

    flow_y = 2.02
    box_w = 1.2
    gap = 0.18
    left_xs = [0.68 + i * (box_w + gap) for i in range(4)]
    left_steps = [
        ("RGB", "几何增强后\n[3,H,W]"),
        ("反向 ISP", "sRGB^-1\ninverse smoothstep"),
        ("随机相机", "CCM / WB\n曝光 / 噪声"),
        ("raw4", "pack Bayer\n[R,Gr,Gb,B]"),
    ]
    for x, (title, body) in zip(left_xs, left_steps):
        add_round_box(slide, x, flow_y, box_w, 0.68, title, body, fill="#FFFFFF", line="#9EC9CF")
    for x in [left_xs[i] + box_w + 0.02 for i in range(3)]:
        add_text(slide, x, flow_y + 0.22, 0.16, 0.18, "→", size=14, color="#3C8994", bold=True, align=PP_ALIGN.CENTER)

    right_xs = [7.22 + i * (box_w + gap) for i in range(4)]
    right_steps = [
        ("RGB", "sRGB 输入\n8-bit / float"),
        ("InvISP", "learned raw-RGB\n或 analytic 近似"),
        ("显式反转", "inverse WB\nfixed RGGB"),
        ("variants", "normal / dark /\nover"),
    ]
    for x, (title, body) in zip(right_xs, right_steps):
        add_round_box(slide, x, flow_y, box_w, 0.68, title, body, fill="#FFFFFF", line="#E3B27F")
    for x in [right_xs[i] + box_w + 0.02 for i in range(3)]:
        add_text(slide, x, flow_y + 0.22, 0.16, 0.18, "→", size=14, color="#C36C1C", bold=True, align=PP_ALIGN.CENTER)

    add_bullets(
        slide,
        0.73,
        2.98,
        5.25,
        1.28,
        [
            "正式配置：raw4 / raw tensor / halfres packed Bayer",
            "训练 randomize=True；KITTI eval 用同一 preset 的 deterministic 版本",
            "sensor_linear_dual 在 ETH3D 与 RobotCar 子 preset 间采样；CFA 可变",
        ],
        color="#2F454C",
        size=8.2,
    )
    add_bullets(
        slide,
        7.27,
        2.98,
        5.25,
        1.28,
        [
            "0524 脚本是离线 NumPy/PIL 工具：保存 npz/json/preview",
            "严格路线依赖外部 InvISP raw-RGB；可运行路线使用 analytic 近似",
            "normal/dark/over 是显式 variant policy，不混用旧 preset 的曝光/噪声范围",
        ],
        color="#5A3B1F",
        size=8.2,
    )

    add_round_box(
        slide,
        0.76,
        4.42,
        5.1,
        0.88,
        "当前方法的含义",
        "把 RGB 合成为“多相机统计驱动”的 pseudo-RAW，适合在线训练增强；但它不是 RAW-Adapter 论文口径的 InvISP。",
        fill="#DDEDEF",
        line="#76A9B1",
        title_color="#12343B",
    )
    add_round_box(
        slide,
        7.3,
        4.42,
        5.1,
        0.88,
        "InvISP 方法的含义",
        "把 RGB 先投影到 raw-RGB，再做 inverse WB 与 RGGB mosaic；更贴近 RAW-Adapter 描述，但必须作为新方法显式接入。",
        fill="#FFE6CE",
        line="#D98A42",
        title_color="#4A2606",
    )

    add_round_box(
        slide,
        0.55,
        6.08,
        12.25,
        0.72,
        "结论 / 下一步",
        "两者输出接口可以统一为 raw4，但实验语义不能混用。正式比较时应新增 unprocessing_method 分支，并在训练与 KITTI eval 配置中显式保存 backend、CFA、gain、variant/noise policy。",
        fill="#FFFFFF",
        line="#B8C2C8",
        title_color="#102027",
        body_color="#2D3A42",
    )

    add_text(
        slide,
        0.46,
        7.05,
        12.45,
        0.2,
        "Sources: foundation/engine/transforms/unprocessing.py; foundation/engine/datasets/vkitti2_raw.py; plans/0524_unprocessing/raw_adapter_style_unprocessing.md; plans/0524_unprocessing/0525_online_raw_adapter_unprocessing_execution_plan.md",
        size=6.5,
        color="#75838C",
    )

    add_ram_to_residual_slide(prs)
    add_network_training_flow_slide(prs)
    add_loss_explanation_slide(prs)
    add_raw_d0_network_detail_slide(prs)
    add_c1_c2_control_comparison_slide(prs)
    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build_deck()
